#AIzaSyAvlTVFIaMFrXoCy1OYDVyODDINJybMqGM


import os
import io
import re
import csv
import json
import time
import traceback
import pandas as pd
import streamlit as st
from fpdf import FPDF
from docx import Document
from dotenv import load_dotenv
from atlassian import Confluence
import google.generativeai as genai
from bs4 import BeautifulSoup
from io import StringIO, BytesIO
import difflib
import warnings
import requests
import re
from sentence_transformers import SentenceTransformer
from PIL import Image, UnidentifiedImageError
from docx.shared import Inches
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import json
# Optional imports for video summarizer
import speech_recognition as sr
from flask import Flask, request, jsonify
from flask_cors import CORS
import tempfile
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import streamlit.components.v1 as components

app = Flask(__name__)
CORS(app)
def record_voice_input():
    recognizer = sr.Recognizer()
    mic = sr.Microphone()

    with mic as source:
        st.info("🎙️ Listening... Speak now.")
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source)

    try:
        text = recognizer.recognize_google(audio)
        st.success(f"📝 Transcribed: {text}")
        return text
    except sr.UnknownValueError:
        st.error("❌ Could not understand audio.")
    except sr.RequestError:
        st.error("❌ Could not connect to the recognition service.")
    
    return ""

try:
    import ffmpeg
except ImportError:
    ffmpeg = None

warnings.filterwarnings("ignore")

# Load environment variables
load_dotenv()

# ------------- Shared Helper Functions -------------
def remove_emojis(text):
    emoji_pattern = re.compile(
        "["
        u"\U0001F600-\U0001F64F"
        u"\U0001F300-\U0001F5FF"
        u"\U0001F680-\U0001F6FF"
        u"\U0001F1E0-\U0001F1FF"
        "]+", flags=re.UNICODE)
    no_emoji = emoji_pattern.sub(r'', text)
    return no_emoji.encode('latin-1', 'ignore').decode('latin-1')


def extract_audio_ffmpeg(video_path, audio_path):
    try:
        ffmpeg.input(video_path).output(audio_path, acodec='mp3', vn=None).run(overwrite_output=True)
    except ffmpeg.Error as e:
        raise RuntimeError(f"FFmpeg error: {e.stderr.decode()}")


def transcribe_with_assemblyai(audio_path):
    api_key = os.getenv("ASSEMBLYAI_API_KEY")
    if not api_key:
        raise ValueError("Missing ASSEMBLYAI_API_KEY in environment.")

    # 🔹 Upload audio file
    with open(audio_path, 'rb') as f:
        upload_response = requests.post(
            "https://api.assemblyai.com/v2/upload",
            headers={"authorization": api_key},
            data=f
        )
    if upload_response.status_code != 200:
        raise RuntimeError("Upload failed: " + upload_response.text)

    audio_url = upload_response.json()["upload_url"]

    # 🔹 Start transcription (✅ no extra fields)
    transcript_response = requests.post(
        "https://api.assemblyai.com/v2/transcript",
        headers={
            "authorization": api_key,
            "content-type": "application/json"
        },
        json={"audio_url": audio_url}
    )

    if transcript_response.status_code != 200:
        raise RuntimeError("Transcription request failed: " + transcript_response.text)

    transcript_id = transcript_response.json()["id"]

    # 🔁 Polling for result
    while True:
        polling_response = requests.get(
            f"https://api.assemblyai.com/v2/transcript/{transcript_id}",
            headers={"authorization": api_key}
        )
        data = polling_response.json()

        if data["status"] == "completed":
            words = data.get("words", [])
            if not words:
                return data["text"]  # fallback to plain if no word timestamps

            # ✅ Build [min:sec] formatted lines
            lines = []
            current_time = ""
            current_line = []
            for word in words:
                start_ms = word.get("start", 0)
                timestamp = f"[{start_ms // 60000}:{(start_ms % 60000) // 1000:02}]"

                if timestamp != current_time:
                    if current_line:
                        lines.append(f"{current_time} {' '.join(current_line)}")
                    current_time = timestamp
                    current_line = [word["text"]]
                else:
                    current_line.append(word["text"])

            if current_line:
                lines.append(f"{current_time} {' '.join(current_line)}")

            return "\n".join(lines)

        elif data["status"] == "error":
            raise RuntimeError("Transcription failed: " + data.get("error", "Unknown error"))
        time.sleep(2)

def clean_html(html_content):
    soup = BeautifulSoup(html_content, "html.parser")
    return soup.get_text(separator="\n")

def create_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for line in text.split('\n'):
        pdf.multi_cell(0, 10, line)
    return io.BytesIO(pdf.output(dest='S').encode('latin1'))

def create_docx(text):
    doc = Document()
    for line in text.split('\n'):
        doc.add_paragraph(line)
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_csv(text):
    output = io.StringIO()
    writer = csv.writer(output)
    for line in text.strip().split('\n'):
        writer.writerow([line])
    return io.BytesIO(output.getvalue().encode())

def create_json(text):
    return io.BytesIO(json.dumps({"response": text}, indent=4).encode())

def create_html(text):
    html = f"<html><body><pre>{text}</pre></body></html>"
    return io.BytesIO(html.encode())

def create_txt(text):
    return io.BytesIO(text.encode())

# ------------- Feature 1: AI Powered Search -------------
@st.cache_resource
def load_faq_data():
    with open("faq_dataset.json", "r", encoding="utf-8") as f:
        data = json.load(f)
    model = SentenceTransformer("all-MiniLM-L6-v2")
    questions = [item["question"] for item in data]
    embeddings = model.encode(questions, normalize_embeddings=True)
    return data, embeddings, model

def get_faq_response(user_query):
    faq_data, faq_embeddings, embed_model = load_faq_data()
    query_embedding = embed_model.encode([user_query], normalize_embeddings=True)
    similarity_scores = cosine_similarity(query_embedding, faq_embeddings)[0]
    best_idx = np.argmax(similarity_scores)
    return faq_data[best_idx]["answer"]



def feature_1():
    st.title("🔗 Confluence AI Powered Search")
    
    # Get query parameters for auto-selection
    query_params = st.query_params
    auto_space_raw = query_params.get("space")
    auto_space = auto_space_raw[0] if isinstance(auto_space_raw, list) else auto_space_raw
    raw_page = query_params.get("page")
    auto_page = raw_page[0] if isinstance(raw_page, list) else raw_page

    @st.cache_resource
    def init_confluence():
        try:
            return Confluence(
                url=os.getenv('CONFLUENCE_BASE_URL'),
                username=os.getenv('CONFLUENCE_USER_EMAIL'),
                password=os.getenv('CONFLUENCE_API_KEY'),
                timeout=10
            )
        except Exception as e:
            st.error(f"Confluence initialization failed: {str(e)}")
            return None

    def init_ai():
        genai.configure(api_key=os.getenv("GENAI_API_KEY"))
        return genai.GenerativeModel("models/gemini-1.5-flash-8b-latest")

    def clean_html(raw_html):
        soup = BeautifulSoup(raw_html, "html.parser")
        return soup.get_text()

    confluence = init_confluence()
    ai_model = init_ai()
    selected_pages = []
    full_context = ""

    if confluence:
        st.success("✅ Connected to Confluence!")

        # Auto-select space if available from query params
        if auto_space:
            space_key = auto_space
            st.success(f"📦 Auto-detected space from URL: {space_key}")
        else:
            try:
                # Fetch spaces
                spaces = confluence.get_all_spaces(start=0, limit=100)["results"]
                space_options = {f"{s['name']} ({s['key']})": s['key'] for s in spaces}
                space_names = list(space_options.keys())

                # Dropdown with no default selection and no fake placeholder in options
                selected_space_label = st.selectbox(
                    label="Select a space key:",
                    options=space_names,
                    index=None,                         # <--- No item selected by default
                    placeholder="Choose a space"       # <--- Shown as grayed-out inside the dropdown
                )

                if selected_space_label is None:
                    st.info("Please select a space to continue.")
                    return

                space_key = space_options[selected_space_label]
            except Exception as e:
                st.error(f"Error fetching spaces: {str(e)}")
                return

        if space_key:
            try:
                # Fetch pages for selected space
                pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=100)
                all_titles = [p["title"] for p in pages]

                select_all = st.checkbox("Select All Pages")
                selected_titles = st.multiselect("Select Page(s):", all_titles, default=all_titles if select_all else [])
                show_content = st.checkbox("Show Page Content")

                selected_pages = [p for p in pages if p["title"] in selected_titles]

                if selected_pages:
                    st.success(f"✅ Loaded {len(selected_pages)} page(s).")
                    for page in selected_pages:
                        page_id = page["id"]
                        page_data = confluence.get_page_by_id(page_id, expand="body.storage")
                        raw_html = page_data["body"]["storage"]["value"]
                        text_content = clean_html(raw_html)
                        full_context += f"\n\nTitle: {page['title']}\n{text_content}"
                        if show_content:
                            with st.expander(f"📄 {page['title']}"):
                                st.markdown(raw_html, unsafe_allow_html=True)
                else:
                    st.warning("Please select at least one page.")
            except Exception as e:
                st.error(f"Error fetching pages: {str(e)}")

    else:
        st.error("❌ Connection to Confluence failed.")
        
    if confluence and selected_pages:
        st.subheader("🤖 Generate AI Response")
        query = st.text_input("Enter your question:")
        if st.button("Generate Answer"):
            if query and full_context:
                try:
                    prompt = (
                        f"Answer the following question using the provided Confluence page content as context.\n"
                        f"Context:\n{full_context}\n\n"
                        f"Question: {query}\n"
                        f"Instructions: Begin with the answer based on the context above. Then, if applicable, supplement with general knowledge."
                    )
                    response = ai_model.generate_content(prompt)
                    st.session_state.ai_response = response.text.strip()
                except Exception as e:
                    st.error(f"AI generation failed: {str(e)}")
            else:
                st.error("Please enter a query.")
                
    if "ai_response" in st.session_state:
        st.markdown("### 💬 AI Response")
        st.markdown(st.session_state.ai_response)
        file_name = st.text_input("Enter file name (without extension):", value="ai_response")
        export_format = st.selectbox("Choose file format to export:", ["TXT", "PDF", "Markdown", "HTML", "DOCX", "CSV", "JSON"])
        export_map = {
            "TXT": (create_txt, "text/plain", ".txt"),
            "PDF": (create_pdf, "application/pdf", ".pdf"),
            "Markdown": (create_txt, "text/markdown", ".md"),
            "HTML": (create_html, "text/html", ".html"),
            "DOCX": (create_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
            "CSV": (create_csv, "text/csv", ".csv"),
            "JSON": (create_json, "application/json", ".json")
        }
        if file_name:
            creator_func, mime, ext = export_map[export_format]
            buffer = creator_func(st.session_state.ai_response)
            st.download_button(
                label="📥 Download File",
                data=buffer,
                file_name=f"{file_name.strip() or 'ai_response'}{ext}",
                mime=mime
            )

        st.markdown("---")
        st.subheader("📝 Save to Confluence Page")

        # Use auto_page as default if available
        if auto_page:
            target_page_title = auto_page
            st.success(f"📄 Auto-selected page to update: {target_page_title}")
        else:
            target_page_title = st.text_input("Enter the Confluence page title to save this to:")

        if st.button("✏️ Save AI Response to Confluence"):
            if target_page_title:
                try:
                    matching_pages = [p for p in selected_pages if p["title"] == target_page_title]
                    if not matching_pages:
                        st.error("Page not found in selected pages.")
                    else:
                        page_id = matching_pages[0]["id"]
                        existing_page = confluence.get_page_by_id(page_id, expand="body.storage")
                        existing_content = existing_page["body"]["storage"]["value"]

                        updated_body = f"{existing_content}<hr/><h3>AI Response</h3><p>{st.session_state.ai_response.replace('\n', '<br>')}</p>"

                        confluence.update_page(
                            page_id=page_id,
                            title=target_page_title,
                            body=updated_body,
                            representation="storage"
                        )
                        st.success("✅ AI response saved to Confluence page.")
                except Exception as e:
                    st.error(f"❌ Failed to update page: {str(e)}")
            else:
                st.warning("Please enter a page title.")

# ------------- Feature 2: Video Summarizer -------------
def feature_2():
    st.title("📄 Confluence Video Summarizer")
    # Get query parameters for auto-selection
    query_params = st.query_params
    auto_space_raw = query_params.get("space")
    auto_space = auto_space_raw[0] if isinstance(auto_space_raw, list) else auto_space_raw
    raw_page = query_params.get("page")
    auto_page = raw_page[0] if isinstance(raw_page, list) else raw_page

    @st.cache_resource
    def init_confluence():
        try:
            return Confluence(
                url=os.getenv("CONFLUENCE_BASE_URL"),
                username=os.getenv("CONFLUENCE_USER_EMAIL"),
                password=os.getenv("CONFLUENCE_API_KEY"),
                timeout=30
            )
        except Exception as e:
            st.error(f"Confluence init failed: {e}")
            return None

    genai.configure(api_key=os.getenv("GENAI_API_KEY"))
    ai_model = genai.GenerativeModel("models/gemini-1.5-flash-8b-latest")
    confluence = init_confluence()

    if ffmpeg is None:
        st.error("Video summarizer dependency `ffmpeg` is not installed. Please run `pip install ffmpeg-python`.")
        return

    if confluence:
        st.success("✅ Connected to Confluence!")

        if auto_space:
            space_key = auto_space
            st.success(f"📦 Auto-detected space from URL: {space_key}")
        else:
            space_key = st.text_input("Enter your Confluence Space Key:")
            
        if space_key:
            try:
                pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=100)
                page_titles = [p["title"] for p in pages]
                selected_pages = st.multiselect("Select Pages to Process:", page_titles)
                if selected_pages:
                    summaries = []
                    for page in pages:
                        if page["title"] not in selected_pages:
                            continue
                        title = page["title"]
                        page_id = page["id"]
                        st.markdown(f"### 🎬 Processing: {title}")
                        attachments = confluence.get(f"/rest/api/content/{page_id}/child/attachment?limit=50")
                        for attachment in attachments["results"]:
                            video_name = attachment["title"].strip()
                            if not video_name.lower().endswith(".mp4"):
                                continue
                            session_key = f"{page_id}_{video_name}".replace(" ", "_")
                            with st.container():
                                if session_key not in st.session_state:
                                    with st.spinner("📥 Downloading and processing..."):
                                        try:
                                            video_url = attachment["_links"]["download"]
                                            full_url = f"{os.getenv('CONFLUENCE_BASE_URL').rstrip('/')}{video_url}"
                                            safe_title = re.sub(r'[<>:"/\\|?*]', '_', title)
                                            safe_video_name = re.sub(r'[<>:"/\\|?*]', '_', video_name)
                                            output_folder = "videos"
                                            os.makedirs(output_folder, exist_ok=True)

                                            # ✅ Sanitize filename to remove illegal Windows characters
                                            safe_title = re.sub(r'[<>:"/\\|?*]', '_', title)
                                            safe_video_name = re.sub(r'[<>:"/\\|?*]', '_', video_name)

                                            # ✅ Final path to save the video
                                            local_path = os.path.join(output_folder, f"{safe_title}_{safe_video_name}".replace(" ", "_"))

                                            # ✅ Download and write video to file
                                            video_data = confluence._session.get(full_url).content
                                            with open(local_path, "wb") as f:
                                                f.write(video_data)

                                            # ✅ Verify the file exists before running FFmpeg
                                            if os.path.exists(local_path):
                                                st.success(f"✅ File successfully saved: {local_path}")
                                            else:
                                                st.error(f"❌ File save failed! Path not found: {local_path}")
                                                continue
                                            # if os.path.exists(local_path):
                                            #     st.success(f"✅ File successfully saved: {local_path}")
                                            # else:
                                            #     st.error(f"❌ File save failed! Path not found: {local_path}")
                                            # continue
                                            extract_audio_ffmpeg(local_path, "temp_audio.mp3")
                                            if os.path.exists("temp_audio.mp3"):
                                                st.success("✅ Audio extracted successfully.")
                                            else:
                                                st.error("❌ Audio extraction failed. FFmpeg did not create temp_audio.mp3")
                                                continue

                                            transcript = transcribe_with_assemblyai("temp_audio.mp3")
    
                                            # Generate summary
                                            quote_prompt = f"Set a title \"Quotes:\" in bold. Extract powerful or interesting quotes:\n{transcript}"
                                            quotes = ai_model.generate_content(quote_prompt).text
                                            summary_prompt = (
                                                f"Start with title as \"Summary:\" in bold, followed by a paragraph.\n"
                                                "**Timestamps:**\n"
                                                "Extract and list only 5–7 important moments from the following transcript.\n"
                                                "Each moment should be one full sentence with the [min:sec] timestamp.\n\n"
                                                f"Transcript:\n{transcript}"
                                            )
                                            summary = ai_model.generate_content(summary_prompt).text
    
                                            st.session_state[session_key] = {
                                                "transcript": transcript,
                                                "summary": summary,
                                                "quotes": quotes
                                            }
                                        except Exception as e:
                                            st.error(f"❌ Error: {e}")
                                            continue
    
                                content = st.session_state[session_key]
                                st.markdown(f"#### 📄 {video_name}")
                                st.markdown(content["quotes"])
                                st.markdown(content["summary"])
                                summaries.append((f"{title}_{video_name}", content["summary"], content["quotes"], content["transcript"]))
    
                                # Q&A Section with form to prevent full rerun
                                q_key = f"{session_key}_question"
                                q_response_key = f"{session_key}_response"
                                q_response_cache_key = f"{session_key}_last_question"
                                
                                # 🎤 Voice input button OUTSIDE the form
                                st.markdown("##### 🎤 Use your voice to ask a question")
                                if st.button("🎤 Speak", key=f"mic_{video_name}"):
                                    spoken_text = record_voice_input()
                                    if spoken_text:
                                        st.session_state[f"ask_{video_name}"] = spoken_text

                                # ✅ Text + Submit inside the form
                                with st.form(key=f"{session_key}_qa_form"):
                                    user_question = st.text_input(
                                        f"💬 Ask a question about {video_name}:",
                                        value=st.session_state.get(f"ask_{video_name}", ""),
                                        key=f"ask_{video_name}_text"
                                    )
                                    submit = st.form_submit_button("🧠 Ask")

                                    if submit and user_question:
                                        if (q_response_cache_key not in st.session_state or
                                            st.session_state[q_response_cache_key] != user_question):
                                            with st.spinner("🤖 Generating answer..."):
                                                try:
                                                    response = ai_model.generate_content(
                                                        f"Answer this in detail based on the video transcription:\n{content['transcript']}\n\nQuestion: {user_question}"
                                                    )
                                                    st.session_state[q_response_key] = response.text
                                                    st.session_state[q_response_cache_key] = user_question
                                                except Exception as e:
                                                    st.error(f"❌ Failed to generate answer: {e}")

                                # Always display last answer if available
                                if q_response_key in st.session_state:
                                    st.markdown(f"**Answer:** {st.session_state[q_response_key]}")


                            # Individual File Download with custom name input
                            file_base = f"{title}_{video_name}".replace(" ", "_").replace(":", "")
                            filename_input_key = f"{session_key}_custom_filename"
                            
                            custom_filename = st.text_input(
                                label="📝 Enter filename (without extension):",
                                value=file_base,
                                key=filename_input_key
                            )
                            format_choice = st.selectbox(
                                f"Choose format for {video_name}:", ["PDF", "TXT"], key=f"{session_key}_format"
                            )
                            file_name = f"{custom_filename.strip() or file_base}.{format_choice.lower()}"
                            export_content = f"{content['quotes']}\n\n{content['summary']}"

                            if format_choice == "PDF":
                                pdf = FPDF()
                                pdf.add_page()
                                pdf.set_font("Arial", size=12)
                                for line in remove_emojis(export_content).split("\n"):
                                    pdf.multi_cell(0, 10, line)
                                file_data = pdf.output(dest="S").encode("latin-1")
                                mime = "application/pdf"
                            else:
                                file_data = export_content.encode("utf-8")
                                mime = "text/plain"

                            st.download_button(
                                label=f"📥 Download {file_name}",
                                data=BytesIO(file_data),
                                file_name=file_name,
                                mime=mime,
                                key=f"{session_key}_download_button"
                            )
                            # ---------------- FAQ BOT SECTION ----------------
                            st.markdown("---")
                            st.subheader("🤖 Ask General Questions (FAQ Bot)")
                            user_faq_query = st.text_input("Type your question about how the summarizer works, downloading, errors, etc.:", key="faq_input")

                            if user_faq_query:
                                with st.spinner("🤖 Searching FAQs..."):
                                    try:
                                        response = get_faq_response(user_faq_query)
                                        st.success("✅ FAQ Bot Answer:")
                                        st.markdown(f"**{response}**")
                                    except Exception as e:
                                        st.error(f"❌ FAQ bot failed: {e}")


                    if len(summaries) > 1:
                        st.markdown("## 📦 Export All Summaries")
                        all_text = ""
                        all_text = ""
                        for t, s, q, _ in summaries:
                            all_text += f"\n\n---\n\n{t}\n\nQuotes:\n{q}\n\nSummary:\n{s}\n"


                        file_name = st.text_input("Filename (without extension):", value="All_Summaries")
                        export_format = st.selectbox("Format:", ["PDF", "TXT"])
                        if export_format == "PDF":
                            pdf = FPDF()
                            pdf.add_page()
                            pdf.set_font("Arial", size=12)
                            for line in remove_emojis(all_text).split("\n"):
                                pdf.multi_cell(0, 10, line)
                            file_data = pdf.output(dest="S").encode("latin-1")
                            mime = "application/pdf"
                            ext = "pdf"
                        else:
                            file_data = all_text.encode("utf-8")
                            mime = "text/plain"
                            ext = "txt"

                        st.download_button(
                            label=f"📥 Download All as {ext.upper()}",
                            data=BytesIO(file_data),
                            file_name=f"{file_name.strip() or 'All_Summaries'}.{ext}",
                            mime=mime
                        )
                        # Save to Confluence functionality
                    if summaries:
                        st.markdown("---")
                        st.subheader("📝 Save to Confluence Page")
                        
                        # Use auto_page as default if available
                        if auto_page:
                            target_page_title = auto_page
                            st.success(f"📄 Auto-selected page to update: {target_page_title}")
                        else:
                            target_page_title = st.text_input("Enter the Confluence page title to save summaries to:")
                        
                        if st.button("✏️ Save Video Summaries to Confluence"):
                            if target_page_title:
                                try:
                                    matching_pages = [p for p in pages if p["title"] == target_page_title]
                                    if not matching_pages:
                                        st.error("Page not found in selected pages.")
                                    else:
                                        page_id = matching_pages[0]["id"]
                                        existing_page = confluence.get_page_by_id(page_id, expand="body.storage")
                                        existing_content = existing_page["body"]["storage"]["value"]
                                        
                                        # Create summary content
                                        summary_content = "<hr/><h3>Video Summaries</h3>"
                                        for summary_title, summary, quotes, _ in summaries:
                                            summary_content += f"<h4>{summary_title}</h4>"
                                            summary_content += f"<h5>Quotes:</h5><p>{quotes.replace(chr(10), '<br>')}</p>"
                                            summary_content += f"<h5>Summary:</h5><p>{summary.replace(chr(10), '<br>')}</p>"
                                            summary_content += "<hr/>"
                                        
                                        updated_body = existing_content + summary_content
                                        
                                        confluence.update_page(
                                            page_id=page_id,
                                            title=target_page_title,
                                            body=updated_body,
                                            representation="storage"
                                        )
                                        st.success("✅ Video summaries saved to Confluence page.")
                                        # ---------------- FAQ BOT FIELD (SEPARATE) ----------------
                                        st.markdown("---")
                                        st.subheader("🤖 FAQ Bot – Ask General Questions")
                                        st.markdown("Ask about how to use this feature, downloading, errors, saving to Confluence, etc.")

                                        user_faq_query = st.text_input("💬 What do you want to ask?", key="faq_input_box")

                                        if user_faq_query:
                                            with st.spinner("🔍 Thinking..."):
                                                try:
                                                    faq_answer = get_faq_response(user_faq_query)
                                                    st.success("✅ FAQ Answer:")
                                                    st.markdown(f"**{faq_answer}**")
                                                except Exception as e:
                                                    st.error(f"❌ FAQ Bot failed: {e}")

                                except Exception as e:
                                    st.error(f"❌ Failed to update page: {str(e)}")
                            else:
                                st.warning("Please enter a page title.")
            except Exception as e:
                st.error(f"Error loading pages: {e}")
    else:
        st.error("❌ Could not connect to Confluence.")

# ------------- Feature 3: Code Assistant -------------
def feature_3():
    import re
    st.title("🔗 Confluence AI Code Assistant")
    
    # Get query parameters for auto-selection
    query_params = st.query_params
    auto_space_raw = query_params.get("space")
    auto_space = auto_space_raw[0] if isinstance(auto_space_raw, list) else auto_space_raw
    raw_page = query_params.get("page")
    auto_page = raw_page[0] if isinstance(raw_page, list) else raw_page
    
    @st.cache_resource
    def init_confluence():
        try:
            return Confluence(
                url=os.getenv('CONFLUENCE_BASE_URL'),
                username=os.getenv('CONFLUENCE_USER_EMAIL'),
                password=os.getenv('CONFLUENCE_API_KEY'),
                timeout=10
            )
        except Exception as e:
            st.error(f"Confluence initialization failed: {str(e)}")
            return None
    def init_ai():
        genai.configure(api_key=os.getenv("GENAI_API_KEY"))
        return genai.GenerativeModel("models/gemini-1.5-flash-8b-latest")
    def strip_code_fences(text: str) -> str:
        return re.sub(r"^```[a-zA-Z]*\n|```$", "", text.strip(), flags=re.MULTILINE)
    def extract_visible_code(html_content: str) -> str:
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup.find_all(['pre', 'code']):
            code_text = tag.get_text()
            if code_text.strip():
                return code_text
        return soup.get_text(separator="\n").strip()
    def detect_language_from_content(content: str) -> str:
        if "<?xml" in content:
            return "xml"
        if "<html" in content.lower() or "<!DOCTYPE html>" in content:
            return "html"
        if content.strip().startswith("{") or content.strip().startswith("["):
            return "json"
        if re.search(r"\bclass\s+\w+", content) and "public" in content:
            return "java"
        if "#include" in content:
            return "cpp"
        if "def " in content:
            return "python"
        if "function" in content or "=>" in content:
            return "javascript"
        return "text"
    def flatten_dict(d, parent_key='', sep='.'):  # for CSV
        items = []
        for k, v in d.items():
            new_key = f"{parent_key}{sep}{k}" if parent_key else k
            if isinstance(v, dict):
                items.extend(flatten_dict(v, new_key, sep=sep).items())
            else:
                items.append((new_key, v))
        return dict(items)
    def create_csv(content):
        try:
            json_data = json.loads(content)
            if isinstance(json_data, dict):
                json_data = [flatten_dict(json_data)]
            if isinstance(json_data, list) and all(isinstance(item, dict) for item in json_data):
                flattened = [flatten_dict(item) for item in json_data]
                fieldnames = sorted(set().union(*(d.keys() for d in flattened)))
                output = io.StringIO()
                writer = csv.DictWriter(output, fieldnames=fieldnames)
                writer.writeheader()
                writer.writerows(flattened)
                return output.getvalue()
            return "Invalid structure for CSV"
        except Exception as e:
            return f"Invalid CSV conversion: {e}"
    def create_txt(content):
        return content
    def create_pdf(content):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Courier", size=10)
        for line in content.splitlines():
            pdf.multi_cell(0, 5, line)
        pdf_output = BytesIO()
        pdf_bytes = pdf.output(dest='S').encode('latin-1')
        pdf_output.write(pdf_bytes)
        pdf_output.seek(0)
        return pdf_output
    def create_docx(content):
        doc = Document()
        for line in content.splitlines():
            doc.add_paragraph(line)
        doc_output = BytesIO()
        doc.save(doc_output)
        doc_output.seek(0)
        return doc_output
    def create_html(content):
        return f"<pre><code>{content}</code></pre>"
    def create_json(content):
        return content
    confluence = init_confluence()
    ai_model = init_ai()
    context = ""
    selected_page = None
    detected_lang = "text"
    if confluence:
        st.success("✅ Connected to Confluence!")
        
        # Auto-select space if available from query params
        if auto_space:
            space_key = auto_space
            st.success(f"📦 Auto-detected space from URL: {space_key}")
        else:
            space_key = st.text_input("Enter your space key:")
            
        if space_key:
            try:
                pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=100)
                page_titles = [p["title"] for p in pages]
                
                # Auto-select page if available from query params
                if auto_page:
                    selected_title = auto_page
                    st.success(f"📄 Auto-selected page: {selected_title}")
                else:
                    selected_title = st.selectbox(
                        "Select a page:",
                        options=page_titles,
                        index=None,
                        placeholder="-- Select a Page --"
                    )
                    
                if selected_title:
                    selected_page = next((p for p in pages if p["title"] == selected_title), None)
                if selected_page:
                    page_id = selected_page["id"]
                    page_content = confluence.get_page_by_id(page_id, expand="body.storage")
                    context = page_content["body"]["storage"]["value"]
                    detected_lang = detect_language_from_content(context)
                    st.success(f"✅ Loaded page: {selected_title}")
                    cleaned_code = extract_visible_code(context)
                    if st.checkbox("📄 Show Page Content"):
                        with st.expander("🔍 Extracted Page Content", expanded=True):
                            st.code(cleaned_code, language=detected_lang)
                    if "summary_response" not in st.session_state:
                        with st.spinner("Generating summary..."):
                            summary_prompt = (
                                f"The following is content (possibly code or structure) from a Confluence page:\n\n{context}\n\n"
                                "Summarize in detailed paragraph"
                            )
                            summary_response = ai_model.generate_content(summary_prompt)
                            st.session_state.summary_response = summary_response.text.strip()
                    st.subheader("📝 Page Summary:")
                    st.markdown(st.session_state.summary_response)
                    st.subheader("🧠 AI Action")
                    ai_action = st.selectbox(
                        "Select action...",
                        ["Select action...", "Summarize Code", "Optimize Performance", "Convert Language", "Generate Documentation", "Refactor Structure", "Security Analysis"]
                    )

                    input_code = st.session_state.get("modified_code", cleaned_code)

                    if ai_action != "Select action..." and st.button("Run AI Action"):
                        action_prompt_map = {
                            "Summarize Code": f"Summarize the following code in clear and concise language:\n\n{input_code}",
                            "Optimize Performance": f"Optimize the following code for performance without changing its functionality:\n\n{input_code}",
                            "Convert Language": f"Convert the following code to another programming language. Suggest a language too:\n\n{input_code}",
                            "Generate Documentation": f"Generate inline documentation and function-level comments for the following code:\n\n{input_code}",
                            "Refactor Structure": f"Refactor the following code to improve structure, readability, and modularity:\n\n{input_code}",
                            "Security Analysis": f"Analyze the following code for security vulnerabilities and suggest improvements:\n\n{input_code}",
                        }

                        with st.spinner(f"Running AI action: {ai_action}"):
                            try:
                                prompt = action_prompt_map.get(ai_action, "")
                                response = ai_model.generate_content(prompt)
                                st.session_state.ai_action_output = strip_code_fences(response.text.strip())
                            except Exception as e:
                                st.error(f"❌ Failed to run action: {e}")
                    if "ai_action_output" in st.session_state:
                        st.subheader(f"📤 Output for: {ai_action}")
                        st.code(st.session_state.ai_action_output, language=detected_lang)

                    st.subheader("✏️ Modify the Code")
                    alter_instruction = st.text_area("Describe the changes you want to make:")
                    if st.button("Modify"):
                        if alter_instruction and cleaned_code:
                            alteration_prompt = (
                                f"The following is a piece of code extracted from a Confluence page:\n\n{cleaned_code}\n\n"
                                f"Please modify this code according to the following instruction:\n'{alter_instruction}'\n\n"
                                "Return the modified code only. No explanation or extra text."
                            )
                            altered_response = ai_model.generate_content(alteration_prompt)
                            st.session_state.modified_code = strip_code_fences(altered_response.text)
                            st.success("✅ Modification Completed")
                    if "modified_code" in st.session_state:
                        st.subheader("🧪 Modified Code Preview")
                        st.code(st.session_state.modified_code, language=detected_lang)
                    st.subheader("🔄 Convert to Another Programming Language")
                    
                    lang_options = [
                                    "Python", "Java", "C#", "JavaScript", "Go", "TypeScript", "C++", "Ruby", "Kotlin",
                                    "Swift", "Rust", "PHP", "Scala", "Perl", "XML", "JSON", "Yang"
                                ]
                    selected_lang = st.selectbox("Select target language:", ["-- Select Language --"] + lang_options)
                    input_code = st.session_state.get("modified_code", cleaned_code)
                    original_lang = detected_lang.lower()
                    target_lang = selected_lang.lower() if selected_lang != "-- Select Language --" else ""
                    if selected_lang != "-- Select Language --" and st.button("Convert Structure"):
                        if original_lang == target_lang:
                            st.error("❌ Cannot convert to the same language.")
                        else:
                            convert_prompt = (
                                f"The following is a code structure or data snippet:\n\n{input_code}\n\n"
                                f"Convert this into equivalent {selected_lang} code. Only show the converted code."
                            )
                            lang_response = ai_model.generate_content(convert_prompt)
                            st.session_state.converted_code = strip_code_fences(lang_response.text)
                    if "converted_code" in st.session_state:
                        st.subheader(f"🔁 Converted to {selected_lang}:")
                        st.code(st.session_state.converted_code, language=selected_lang.lower())
                        file_name = st.text_input("Enter file name (without extension):", value="ai_response")
                        export_format = st.selectbox("Choose file format:", ["TXT", "PDF", "Markdown", "HTML", "DOCX", "CSV", "JSON"])
                        export_map = {
                            "TXT": (create_txt, "text/plain", ".txt"),
                            "PDF": (create_pdf, "application/pdf", ".pdf"),
                            "Markdown": (create_txt, "text/markdown", ".md"),
                            "HTML": (create_html, "text/html", ".html"),
                            "DOCX": (create_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
                            "CSV": (create_csv, "text/csv", ".csv"),
                            "JSON": (create_json, "application/json", ".json")
                        }
                        if file_name:
                            creator_func, mime, ext = export_map[export_format]
                            buffer = creator_func(st.session_state.converted_code)
                            st.download_button(
                                label="📥 Download File",
                                data=buffer,
                                file_name=f"{file_name.strip() or 'ai_response'}{ext}",
                                mime=mime
                            )
                        
                        # Save to Confluence functionality
                        st.markdown("---")
                        st.subheader("📝 Save to Confluence Page")
                        
                        # Use auto_page as default if available, otherwise use current page
                        if auto_page:
                            target_page_title = auto_page
                            st.success(f"📄 Auto-selected page to update: {target_page_title}")
                        else:
                            target_page_title = st.text_input("Enter the Confluence page title to save converted code to:", value=selected_title)
                        
                        if st.button("✏️ Save Converted Code to Confluence"):
                            if target_page_title:
                                try:
                                    matching_pages = [p for p in pages if p["title"] == target_page_title]
                                    if not matching_pages:
                                        st.error("Page not found in selected pages.")
                                    else:
                                        page_id = matching_pages[0]["id"]
                                        existing_page = confluence.get_page_by_id(page_id, expand="body.storage")
                                        existing_content = existing_page["body"]["storage"]["value"]
                                        
                                        # Create code content with proper formatting
                                        code_content = f"<hr/><h3>Converted Code ({selected_lang})</h3>"
                                        code_content += f"<ac:structured-macro ac:name=\"code\"><ac:parameter ac:name=\"language\">{selected_lang.lower()}</ac:parameter>"
                                        code_content += f"<ac:plain-text-body><![CDATA[{st.session_state.converted_code}]]></ac:plain-text-body></ac:structured-macro>"
                                        
                                        updated_body = existing_content + code_content
                                        
                                        confluence.update_page(
                                            page_id=page_id,
                                            title=target_page_title,
                                            body=updated_body,
                                            representation="storage"
                                        )
                                        st.success("✅ Converted code saved to Confluence page.")
                                except Exception as e:
                                    st.error(f"❌ Failed to update page: {str(e)}")
                            else:
                                st.warning("Please enter a page title.")
            except Exception as e:
                st.error(f"Error fetching pages: {str(e)}")
    else:
        st.error("❌ Connection to Confluence failed.")

# ------------- Feature 4: Impact Analyzer -------------
def feature_4():
    import re
    import regex
    def remove_emojis(text):
    # Remove symbols, emojis, and non-latin characters safely
        return regex.sub(r'[\p{So}\p{Sk}\p{Cn}]+', '', text)

    st.title("🧠 Confluence AI Impact Analyzer")
    
    # Get query parameters for auto-selection
    query_params = st.query_params
    auto_space_raw = query_params.get("space")
    auto_space = auto_space_raw[0] if isinstance(auto_space_raw, list) else auto_space_raw
    raw_page = query_params.get("page")
    auto_page = raw_page[0] if isinstance(raw_page, list) else raw_page
    
    @st.cache_resource
    def init_confluence():
        try:
            return Confluence(
                url=os.getenv('CONFLUENCE_BASE_URL'),
                username=os.getenv('CONFLUENCE_USER_EMAIL'),
                password=os.getenv('CONFLUENCE_API_KEY'),
                timeout=10
            )
        except Exception as e:
            st.error(f"Confluence init failed: {str(e)}")
            return None
    genai.configure(api_key=os.getenv("GENAI_API_KEY"))
    model = genai.GenerativeModel("models/gemini-1.5-flash-8b-latest")
    MAX_CHARS = 10000
    def extract_code_blocks(content):
        soup = BeautifulSoup(content, 'html.parser')
        blocks = soup.find_all('ac:structured-macro', {'ac:name': 'code'})
        return '\n'.join(
            block.find('ac:plain-text-body').text
            for block in blocks if block.find('ac:plain-text-body')
        )
    def clean_and_truncate_prompt(text, max_chars=MAX_CHARS):
        text = re.sub(r'<[^>]+>', '', text)
        text = re.sub(r'[^\x00-\x7F]+', '', text)
        return text[:max_chars]
    def safe_generate(prompt, retries=3):
        prompt = clean_and_truncate_prompt(prompt)
        fallback_prompt = "Explain this code change or answer a general question about code quality."
        for i in range(retries):
            try:
                return model.generate_content(prompt).text.strip()
            except Exception as e:
                st.warning(f"Retry {i+1} failed: {e}")
                time.sleep(2)
        st.warning("⚠️ Using fallback response due to repeated errors.")
        return model.generate_content(fallback_prompt).text.strip()
    confluence = init_confluence()
    if confluence:
        st.success("✅ Connected to Confluence")
        
        # Auto-select space if available from query params
        if auto_space:
            space_key = auto_space
            st.success(f"📦 Auto-detected space from URL: {space_key}")
        else:
            space_key = st.text_input("Enter your Confluence Space Key:")
            
        page_titles = []
        if space_key:
            try:
                pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=100)
                page_titles = [p["title"] for p in pages]
            except Exception as e:
                st.error(f"Error fetching pages from space '{space_key}': {e}")
        if page_titles:
            old_page_title = st.selectbox("OLD version code page", options=page_titles, index=None, placeholder="Select a page", key="old_page")
            new_page_title = st.selectbox("NEW version code page", options=page_titles, index=None, placeholder="Select a page", key="new_page")
        else:
            old_page_title = ""
            new_page_title = ""
        if old_page_title and new_page_title:
            try:
                old_page = next((p for p in pages if p["title"] == old_page_title), None)
                new_page = next((p for p in pages if p["title"] == new_page_title), None)
                if old_page and new_page:
                    old_raw = confluence.get_page_by_id(old_page["id"], expand="body.storage")["body"]["storage"]["value"]
                    new_raw = confluence.get_page_by_id(new_page["id"], expand="body.storage")["body"]["storage"]["value"]
                    old_code = extract_code_blocks(old_raw)
                    new_code = extract_code_blocks(new_raw)
                    st.subheader(f"📄 {old_page_title} Code")
                    st.code(old_code or "No code found", language='python')
                    st.subheader(f"📄 {new_page_title} Code")
                    st.code(new_code or "No code found", language='python')
                    if old_code and new_code:
                        old_lines = old_code.splitlines()
                        new_lines = new_code.splitlines()
                        diff = difflib.unified_diff(old_lines, new_lines, fromfile=old_page_title, tofile=new_page_title, lineterm='')
                        full_diff_text = '\n'.join(diff)
                        safe_diff = clean_and_truncate_prompt(full_diff_text)
                        lines_added = sum(1 for l in full_diff_text.splitlines() if l.startswith('+') and not l.startswith('+++'))
                        lines_removed = sum(1 for l in full_diff_text.splitlines() if l.startswith('-') and not l.startswith('---'))
                        total_lines = len(old_lines) or 1
                        percent_change = round(((lines_added + lines_removed) / total_lines) * 100, 2)
                        code_blocks_changed = abs(old_code.count('\n') // 5 - new_code.count('\n') // 5)
                        st.subheader("📈 Change Metrics Dashboard")
                        st.markdown(f"""
                        <div style="border:1px solid #ddd; padding:10px; border-radius:10px; background:#f9f9f9">
                            <ul>
                                <li><b>Lines Added:</b> {lines_added}</li>
                                <li><b>Lines Removed:</b> {lines_removed}</li>
                                <li><b>Percentage Changed:</b> {percent_change}%</li>
                                <li><b>Code Blocks Changed:</b> {code_blocks_changed}</li>
                            </ul>
                        </div>
                        """, unsafe_allow_html=True)
                        if "impact_text" not in st.session_state:
                            st.session_state.impact_text = safe_generate(
                                        f"""Write 2 paragraphs summarizing the overall impact of the following code diff.
                            
                                        Cover only:
                                        - What was changed
                                        - Which parts of the system are affected
                                        - Why this matters
                                        
                                        Keep it within 20 sentences.
                                        
                                        Code Diff:
                                        {safe_diff}"""
                                            )
                        if "rec_text" not in st.session_state:
                                st.session_state.rec_text = safe_generate(
                                        f"""As a senior engineer, write 2 paragraphs suggesting improvements for the following code diff.

                                        Focus on:
                                        - Code quality
                                        - Maintainability
                                        - Any possible optimizations
                                        
                                        Limit to 20 sentences.
                                        
                                        Code Diff:
                                        {safe_diff}"""
                                            )
                        if "risk_text" not in st.session_state:
                            raw_risk = safe_generate(f"Assess the risk of each change in this code diff with severity tags (Low, Medium, High):\n\n{safe_diff}")
                            st.session_state.risk_text = re.sub(
                                r'\b(Low|Medium|High)\b',
                                lambda m: {
                                    'Low': '🟢 Low',
                                    'Medium': '🟡 Medium',
                                    'High': '🔴 High'
                                }[m.group(0)],
                                raw_risk
                            )
                        st.subheader("📌 Impact Analysis Summary")
                        st.markdown(st.session_state.impact_text)
                        st.subheader("✨ AI-Powered Change Recommendations")
                        st.markdown(st.session_state.rec_text)
                        st.subheader("🛡️ Risk Analysis with Severity Levels")
                        st.markdown(st.session_state.risk_text)
                        st.markdown("---")
                        st.header("💬 Ask a Question about the AI Analysis")
                        if "user_question" not in st.session_state:
                            st.session_state.user_question = ""
                        if "qa_answer" not in st.session_state:
                            st.session_state.qa_answer = ""
                        user_question_input = st.text_input("Ask a question about the AI-generated results:")
                        if user_question_input and user_question_input != st.session_state.user_question:
                            st.session_state.user_question = user_question_input
                            context = (
                                f"Summary: {st.session_state.impact_text[:1000]}\n"
                                f"Recommendations: {st.session_state.rec_text[:1000]}\n"
                                f"Risks: {st.session_state.risk_text[:1000]}\n"
                                f"Changes: +{lines_added}, -{lines_removed}, ~{percent_change}%"
                            )
                            qa_prompt = f"""You are an expert AI assistant. Based on the report below, answer the user's question clearly.

{context}

Question: {user_question_input}

Answer:"""
                            st.session_state.qa_answer = safe_generate(qa_prompt)
                        if st.session_state.qa_answer:
                            st.subheader("🤖 AI Answer")
                            st.markdown(st.session_state.qa_answer)
                        st.markdown("---")
                        st.header("📁 Download:")
                        file_name = st.text_input("Enter file name (without extension):", value=f"{new_page_title}_impact")
                        export_format = st.selectbox("Choose file format to export:", ["Markdown (.md)", "PDF (.pdf)", "Text (.txt)"])
                        md_content = f"""# Impact Summary

{st.session_state.impact_text}

## Change Recommendations

{st.session_state.rec_text}

## Risk Analysis

{st.session_state.risk_text}
"""
                        if export_format.startswith("Markdown"):
                            st.download_button(
                                label="📥 Download Markdown",
                                data=md_content.encode("utf-8"),
                                file_name=f"{file_name}.md",
                                mime="text/markdown"
                            )
                        elif export_format.startswith("PDF"):
                            pdf = FPDF()
                            pdf.add_page()
                            pdf.set_font("Arial", size=12)

                            clean_report = remove_emojis(md_content)  # <-- clean content before PDF
                            for line in clean_report.split("\n"):
                                try:
                                    pdf.multi_cell(0, 10, line)
                                except Exception:
                                    # Fallback for any unexpected characters
                                    safe_line = line.encode("latin-1", "replace").decode("latin-1")
                                    pdf.multi_cell(0, 10, safe_line)

                            pdf_bytes = pdf.output(dest='S').encode("latin-1")
                            st.download_button(
                                label="📥 Download PDF",
                                data=BytesIO(pdf_bytes),
                                file_name=f"{file_name}.pdf",
                                mime="application/pdf"
                            )
                        else:
                            st.download_button(
                                label="📥 Download TXT",
                                data=md_content.encode("utf-8"),
                                file_name=f"{file_name}.txt",
                                mime="text/plain"
                            )
                        
                        # Save to Confluence functionality
                        st.markdown("---")
                        st.subheader("📝 Save to Confluence Page")
                        
                        # Use auto_page as default if available, otherwise use new page
                        if auto_page:
                            target_page_title = auto_page
                            st.success(f"📄 Auto-selected page to update: {target_page_title}")
                        else:
                            target_page_title = st.text_input("Enter the Confluence page title to save impact analysis to:", value=new_page_title)
                        
                        if st.button("✏️ Save Impact Analysis to Confluence"):
                            if target_page_title:
                                try:
                                    matching_pages = [p for p in pages if p["title"] == target_page_title]
                                    if not matching_pages:
                                        st.error("Page not found in selected pages.")
                                    else:
                                        page_id = matching_pages[0]["id"]
                                        existing_page = confluence.get_page_by_id(page_id, expand="body.storage")
                                        existing_content = existing_page["body"]["storage"]["value"]
                                        
                                        # Create impact analysis content
                                        impact_content = "<hr/><h3>Impact Analysis Report</h3>"
                                        impact_content += f"<h4>Change Metrics</h4>"
                                        impact_content += f"<ul><li>Lines Added: {lines_added}</li>"
                                        impact_content += f"<li>Lines Removed: {lines_removed}</li>"
                                        impact_content += f"<li>Percentage Changed: {percent_change}%</li>"
                                        impact_content += f"<li>Code Blocks Changed: {code_blocks_changed}</li></ul>"
                                        impact_content += f"<h4>Impact Summary</h4><p>{st.session_state.impact_text.replace(chr(10), '<br>')}</p>"
                                        impact_content += f"<h4>Recommendations</h4><p>{st.session_state.rec_text.replace(chr(10), '<br>')}</p>"
                                        impact_content += f"<h4>Risk Analysis</h4><p>{st.session_state.risk_text.replace(chr(10), '<br>')}</p>"
                                        
                                        updated_body = existing_content + impact_content
                                        
                                        confluence.update_page(
                                            page_id=page_id,
                                            title=target_page_title,
                                            body=updated_body,
                                            representation="storage"
                                        )
                                        st.success("✅ Impact analysis saved to Confluence page.")
                                except Exception as e:
                                    st.error(f"❌ Failed to update page: {str(e)}")
                            else:
                                st.warning("Please enter a page title.")
            except Exception as e:
                st.error(f"Error: {e}")
    else:
        st.error("❌ Connection to Confluence failed.")

# ------------- Feature 5: Test Support Tool -------------
def feature_5():
    st.title("🤖 Confluence AI Test Support Tool")
    
    # Get query parameters for auto-selection
    query_params = st.query_params
    auto_space_raw = query_params.get("space")
    auto_space = auto_space_raw[0] if isinstance(auto_space_raw, list) else auto_space_raw
    raw_page = query_params.get("page")
    auto_page = raw_page[0] if isinstance(raw_page, list) else raw_page
    
    @st.cache_resource
    def init_ai():
        genai.configure(api_key=os.getenv("GENAI_API_KEY"))
        return genai.GenerativeModel("models/gemini-1.5-flash-8b-latest")
    @st.cache_resource
    def init_confluence():
        try:
            return Confluence(
                url=os.getenv('CONFLUENCE_BASE_URL'),
                username=os.getenv('CONFLUENCE_USER_EMAIL'),
                password=os.getenv('CONFLUENCE_API_KEY'),
                timeout=10
            )
        except Exception as e:
            st.error(f"Confluence initialization failed: {str(e)}")
            return None
    ai_model = init_ai()
    confluence = init_confluence()
    if 'strategy_text' not in st.session_state:
        st.session_state.strategy_text = ""
    if 'cross_text' not in st.session_state:
        st.session_state.cross_text = ""
    if 'sensitivity_text' not in st.session_state:
        st.session_state.sensitivity_text = ""
    if 'ai_response' not in st.session_state:
        st.session_state.ai_response = ""
    if confluence:
        st.success("✅ Connected to Confluence!")
        
        # Auto-select space if available from query params
        if auto_space:
            space_key = auto_space
            st.success(f"📦 Auto-detected space from URL: {space_key}")
        else:
            space_key = st.text_input("Enter your Confluence space key:")
            
        if space_key:
            try:
                pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=50)
                titles = [page['title'] for page in pages]
                
                # Auto-select pages if available from query params
            
                selected_code_title = st.selectbox("Select Code Page", options=titles, index=None, placeholder="Choose a code page")
                selected_test_input_title = st.selectbox("Select Test Input Page", options=titles, index=None, placeholder="Choose a test input page")
                    
                code_page = next((p for p in pages if p["title"] == selected_code_title), None)
                test_input_page = next((p for p in pages if p["title"] == selected_test_input_title), None)
                if code_page:
                    code_data = confluence.get_page_by_id(code_page["id"], expand="body.storage")
                    code_content = code_data["body"]["storage"]["value"]
                    st.markdown("### 📘 Confluence Test Strategy Generator")
                    if not st.session_state.strategy_text:
                        with st.spinner("🧪 Generating test strategy..."):
                            prompt_strategy = f"""The following is a code snippet:\n\n{code_content}\n\nBased on this, please generate appropriate test strategies and test cases. Mention types of testing (unit, integration, regression), areas that require special attention, and possible edge cases."""
                            response_strategy = ai_model.generate_content(prompt_strategy)
                            st.session_state.strategy_text = response_strategy.text.strip()
                    st.subheader("📋 Suggested Test Strategies and Test Cases")
                    st.markdown(st.session_state.strategy_text)
                    st.markdown("### 🌐 Cross-Platform Testing Intelligence")
                    if not st.session_state.cross_text:
                        with st.spinner("🧠 Analyzing for cross-platform compatibility..."):
                            prompt_cross_platform = f"""You are a cross-platform UI testing expert. Analyze the following frontend code and generate test strategies. Code:\n{code_content}\n\nInclude: - Desktop, Mobile Web, Tablet test cases - UI/viewport issues in one paragraph - Framework/tool suggestions in one paragraph"""
                            response_cross_platform = ai_model.generate_content(prompt_cross_platform)
                            st.session_state.cross_text = response_cross_platform.text.strip()
                    st.subheader("📋 Suggested Strategy and Test Cases")
                    st.markdown(st.session_state.cross_text)
                if test_input_page:
                    test_data = confluence.get_page_by_id(test_input_page["id"], expand="body.storage")
                    test_input_content = test_data["body"]["storage"]["value"]
                    st.markdown("### 🔒 Data Sensitivity Classifier for Test Inputs")
                    st.code(test_input_content, language="json")
                    if st.button("🔍 Classify Sensitive Data"):
                        with st.spinner("🔐 Analyzing for sensitive data..."):
                            prompt_sensitivity = f"""You are a data privacy expert. Classify sensitive fields (PII, credentials, financial) and provide masking suggestions.\n\nData:\n{test_input_content}"""
                            response_sensitivity = ai_model.generate_content(prompt_sensitivity)
                            st.session_state.sensitivity_text = response_sensitivity.text.strip()
                    if st.session_state.sensitivity_text:
                        st.subheader("📋 Sensitivity Analysis and Recommendations")
                        st.markdown(st.session_state.sensitivity_text)
                if all([
                    st.session_state.strategy_text,
                    st.session_state.cross_text,
                    st.session_state.sensitivity_text
                ]):
                    st.markdown("### 📥 Download Full Report")
                    filename_input = st.text_input("Enter filename (without extension):", value="ai_testing_report", key="filename_input")
                    file_format = st.selectbox("Select file format:", ["TXT", "PDF"], key="format_selector")
                    full_report = (
                        "📘 Test Strategy:\n" + st.session_state.strategy_text + "\n\n"
                        "🌐 Cross-Platform Testing:\n" + st.session_state.cross_text + "\n\n"
                        "🔒 Sensitivity Analysis:\n" + st.session_state.sensitivity_text
                    )
                    filename = f"{filename_input}.{file_format.lower()}"
                    if file_format == "TXT":
                        file_bytes = full_report.encode("utf-8")
                        mime = "text/plain"
                    else:
                        pdf = FPDF()
                        pdf.add_page()
                        pdf.set_auto_page_break(auto=True, margin=15)
                        pdf.set_font("Arial", size=12)
                        clean_report = remove_emojis(full_report)
                        for line in clean_report.split("\n"):
                            try:
                                pdf.multi_cell(0, 10, line)
                            except:
                                pdf.multi_cell(0, 10, line.encode('latin-1', 'replace').decode('latin-1'))
                        pdf_output = pdf.output(dest='S').encode("latin-1")
                        file_bytes = BytesIO(pdf_output).getvalue()
                        mime = "application/pdf"
                    st.download_button(
                        label="📄 Generate and Download File",
                        data=file_bytes,
                        file_name=filename,
                        mime=mime
                    )
                    st.markdown("### 🤖 Ask Questions")
                    user_question = st.text_input("Ask a question about the generated results:")
                    
                    if user_question:
                        with st.spinner("🤖 Thinking..."):
                            prompt_chat = f"""Based on the following content:\n📘 Test Strategy:\n{st.session_state.strategy_text}\n🌐 Cross-Platform Testing:\n{st.session_state.cross_text}\n🔒 Sensitivity Analysis:\n{st.session_state.sensitivity_text}\n\nAnswer this user query: \"{user_question}\" """
                            ai_response = ai_model.generate_content(prompt_chat)
                            st.session_state.ai_response = ai_response.text.strip()
                    else:
                        st.session_state.ai_response = ""  # ❗ Clear previous response if no new question is asked
                    
                    if st.session_state.get("ai_response"):
                        st.markdown(f"**🤖 AI Response:** {st.session_state.ai_response}")
                    
                    # Save to Confluence functionality
                    st.markdown("---")
                    st.subheader("📝 Save to Confluence Page")
                    
                    # Use auto_page as default if available, otherwise use code page
                    if auto_page:
                        target_page_title = auto_page
                        st.success(f"📄 Auto-selected page to update: {target_page_title}")
                    else:
                        target_page_title = st.text_input("Enter the Confluence page title to save test analysis to:", value=selected_code_title if selected_code_title else "")
                    
                    if st.button("✏️ Save Test Analysis to Confluence"):
                        if target_page_title:
                            try:
                                matching_pages = [p for p in pages if p["title"] == target_page_title]
                                if not matching_pages:
                                    st.error("Page not found in selected pages.")
                                else:
                                    page_id = matching_pages[0]["id"]
                                    existing_page = confluence.get_page_by_id(page_id, expand="body.storage")
                                    existing_content = existing_page["body"]["storage"]["value"]
                                    
                                    # Create test analysis content
                                    test_content = "<hr/><h3>Test Analysis Report</h3>"
                                    test_content += f"<h4>Test Strategy</h4><p>{st.session_state.strategy_text.replace(chr(10), '<br>')}</p>"
                                    test_content += f"<h4>Cross-Platform Testing</h4><p>{st.session_state.cross_text.replace(chr(10), '<br>')}</p>"
                                    test_content += f"<h4>Sensitivity Analysis</h4><p>{st.session_state.sensitivity_text.replace(chr(10), '<br>')}</p>"
                                    
                                    if st.session_state.ai_response:
                                        test_content += f"<h4>AI Q&A</h4><p><strong>Question:</strong> {user_question}</p>"
                                        test_content += f"<p><strong>Answer:</strong> {st.session_state.ai_response.replace(chr(10), '<br>')}</p>"
                                    
                                    updated_body = existing_content + test_content
                                    
                                    confluence.update_page(
                                        page_id=page_id,
                                        title=target_page_title,
                                        body=updated_body,
                                        representation="storage"
                                    )
                                    st.success("✅ Test analysis saved to Confluence page.")
                            except Exception as e:
                                st.error(f"❌ Failed to update page: {str(e)}")
                        else:
                            st.warning("Please enter a page title.")
            except Exception as e:
                st.error(f"Error retrieving Confluence data: {str(e)}")
    else:
        st.error("❌ Could not connect to Confluence.")

# ------------- Feature 6: Image Generation and Graph Gen-------------
def feature_6():
    st.title("🖼️ Confluence AI Image Generation and Graph Generation")
    def init_confluence():
        try:
            return Confluence(
                url=os.getenv('CONFLUENCE_BASE_URL'),
                username=os.getenv('CONFLUENCE_USER_EMAIL'),
                password=os.getenv('CONFLUENCE_API_KEY'),
                timeout=10
            )
        except Exception as e:
            st.error(f"Confluence initialization failed: {str(e)}")
            return None

    def init_ai():
        genai.configure(api_key=os.getenv("GENAI_API_KEY"))
        return genai.GenerativeModel("models/gemini-1.5-flash")

    def download_image_bytes(url: str, auth):
        response = requests.get(url, auth=auth)
        if response.status_code == 200:
            return response.content
        else:
            raise Exception(f"Failed to fetch image (Status: {response.status_code})")

    def generate_pdf(image_bytes, summary):
        pdf = FPDF()
        pdf.add_page()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            img = Image.open(BytesIO(image_bytes))
            img.save(tmp_img, format="PNG")
            tmp_img_path = tmp_img.name
        pdf.image(tmp_img_path, x=10, w=180)
        pdf.ln(5)
        safe_summary = summary.encode('latin-1', 'replace').decode('latin-1')
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, safe_summary)
        pdf_output = pdf.output(dest='S').encode('latin-1')
        return BytesIO(pdf_output)

    def generate_docx(image_bytes, summary):
        doc = Document()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            img = Image.open(BytesIO(image_bytes))
            img.save(tmp_img, format="PNG")
            tmp_img_path = tmp_img.name
        doc.add_picture(tmp_img_path, width=Inches(5.5))
        doc.add_paragraph(summary)
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    def generate_txt(summary):
        output = BytesIO()
        output.write(summary.encode())
        output.seek(0)
        return output

    def generate_md(image_bytes, summary):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            img = Image.open(BytesIO(image_bytes))
            img.save(tmp_img, format="PNG")
            image_path = tmp_img.name
        content = f"![Image]({image_path})\n\n{summary}"
        output = BytesIO()
        output.write(content.encode())
        output.seek(0)
        return output

    def plot_grouped_bar(df):
        melted = df.melt(id_vars=[df.columns[0]], var_name="Group", value_name="Count")
        plt.figure(figsize=(10, 6))
        sns.barplot(data=melted, x=melted.columns[0], y="Count", hue="Group")
        plt.xticks(rotation=45)
        plt.title("Grouped Bar Chart")
        plt.tight_layout()
        return plt.gcf()

    def plot_stacked_bar(df):
        df_plot = df.set_index(df.columns[0])
        plt.figure(figsize=(10, 6))
        df_plot.drop(columns="Total", errors="ignore").plot(kind='bar', stacked=True)
        plt.title("Stacked Bar Chart")
        plt.xticks(rotation=45)
        plt.ylabel("Count")
        plt.tight_layout()
        return plt.gcf()

    def plot_line(df):
        df_plot = df.set_index(df.columns[0])
        plt.figure(figsize=(10, 6))
        df_plot.drop(columns="Total", errors="ignore").plot(marker='o')
        plt.title("Line Chart")
        plt.xticks(rotation=45)
        plt.ylabel("Count")
        plt.tight_layout()
        return plt.gcf()

    def plot_pie(df):
        plt.figure(figsize=(7, 6))
        label_col = df.columns[0]
        
        if "Total" in df.columns:
            data = df["Total"]
        else:
            # Fallback: Sum across all numeric columns (except label)
            data = df.iloc[:, 1:].sum(axis=1)

        plt.pie(data, labels=df[label_col], autopct="%1.1f%%", startangle=140)
        plt.title("Pie Chart (Total Responses)")
        plt.tight_layout()
        return plt.gcf()


    def get_image_bytes(fig, fmt):
        buf = BytesIO()
        fig.savefig(buf, format=fmt, bbox_inches="tight")
        buf.seek(0)
        return buf

    def clean_ai_csv(raw_text):
        lines = raw_text.strip().splitlines()
        clean_lines = [
            line.strip() for line in lines
            if ',' in line and not line.strip().startswith("```") and not line.lower().startswith("here")
        ]
        # Remove duplicate headers or malformed lines
        header = clean_lines[0].split(",")
        cleaned_data = [clean_lines[0]]
        for line in clean_lines[1:]:
            if line.split(",")[0] != header[0]:  # skip repeat headers
                cleaned_data.append(line)
        return "\n".join(cleaned_data)

    def generate_chart_pdf(image_bytes):
        pdf = FPDF()
        pdf.add_page()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            tmp_img.write(image_bytes.read())
            tmp_img.flush()
            pdf.image(tmp_img.name, x=10, w=180)
        output = BytesIO()
        pdf.output(output)
        output.seek(0)
        return output

    def generate_chart_docx(image_bytes):
        doc = Document()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            tmp_img.write(image_bytes.read())
            tmp_img.flush()
            doc.add_picture(tmp_img.name, width=Inches(5.5))
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    def generate_chart_pptx(image_bytes):
        prs = Presentation()
        blank_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            tmp_img.write(image_bytes.read())
            tmp_img.flush()
            slide_width = prs.slide_width
            top = Inches(1)
            left = Inches(0.5)
            pic = blank_slide.shapes.add_picture(tmp_img.name, left, top, width=slide_width - Inches(1))
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
    
    confluence = init_confluence()
    ai_model = init_ai()

    if confluence:
        st.success("✅ Connected to Confluence")

        space_key = st.text_input("Enter your space key:")

        if space_key:
            try:
                pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=100)
                all_titles = [p["title"] for p in pages]
                selected_titles = st.multiselect("Select page titles:", options=all_titles)

                for page_title in selected_titles:
                    selected_page = next((p for p in pages if p["title"].strip().lower() == page_title.strip().lower()), None)

                    if selected_page:
                        page_id = selected_page["id"]
                        html_content = confluence.get_page_by_id(page_id=page_id, expand="body.export_view")["body"]["export_view"]["value"]
                        soup = BeautifulSoup(html_content, "html.parser")
                        base_url = os.getenv("CONFLUENCE_BASE_URL")

                        image_urls = list({
                            base_url + img["src"] if img["src"].startswith("/") else img["src"]
                            for img in soup.find_all("img") if img.get("src")
                        })

                        if not image_urls:
                            st.warning(f"⚠️ No embedded images found in the page '{page_title}'.")
                        else:
                            st.subheader(f"📸 Images in: {page_title}")

                            for idx, url in enumerate(image_urls):
                                summary_key = f"summary_{page_title}_{idx}"
                                ready_key = f"ready_{page_title}_{idx}"
                                image_key = f"image_{page_title}_{idx}"
                                graph_df_key = f"graph_df_{page_title}_{idx}"
                                graph_fig_key = f"graph_fig_{page_title}_{idx}"
                                chart_type_key = f"chart_type_{page_title}_{idx}"

                                try:
                                    image_bytes = download_image_bytes(url, auth=(os.getenv('CONFLUENCE_USER_EMAIL'), os.getenv('CONFLUENCE_API_KEY')))
                                    st.session_state[image_key] = image_bytes

                                    try:
                                        image_pil = Image.open(BytesIO(image_bytes))
                                        st.image(image_pil, use_container_width=True)
                                    except UnidentifiedImageError:
                                        st.warning(f"Image {idx} could not be loaded.")
                                        continue

                                    if st.button("Summarize", key=f"summarize_{page_title}_{idx}"):
                                        with st.spinner("Generating summary..."):
                                            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                                                tmp.write(image_bytes)
                                                tmp.flush()
                                                uploaded = genai.upload_file(
                                                    path=tmp.name,
                                                    mime_type="image/png",
                                                    display_name=f"confluence_image_{page_title}_{idx}.png"
                                                )

                                            prompt = (
                                                "You are analyzing a technical image from a documentation page. "
                                                "If it's a chart or graph, explain what is shown in detail. "
                                                "If it's code, summarize what the code does. "
                                                "Avoid mentioning filenames or metadata. Provide an informative analysis in 2 paragraphs."
                                            )

                                            try:
                                                response = ai_model.generate_content([uploaded, prompt])
                                                st.session_state[summary_key] = response.text.strip()
                                                st.session_state[ready_key] = True           

                                            except Exception as e:
                                                st.error(f"AI failed to generate summary: {e}")

                                    if st.session_state.get(ready_key):
                                        anchor_id = f"summary_{page_title}_{idx}"
                                                                                            
                                        st.subheader("Gemini Summary")
                                        summary = st.session_state[summary_key]
                                        st.success(summary)

                
                                        st.markdown("Generate AI response")
                                        ai_response_key = f"ai_response_{page_title}_{idx}"
                                        user_question_key = f"user_question_{page_title}_{idx}"

                                        # Persist user input
                                        if user_question_key not in st.session_state:
                                            st.session_state[user_question_key] = ""

                                        user_question = st.text_input(
                                            "Enter the question:",
                                            key=user_question_key,
                                            value=st.session_state[user_question_key],
                                            placeholder="Ask the Question?"
                                        )

                                        if user_question and st.session_state.get(ai_response_key) is None:
                                            with st.spinner("Generating AI response..."):
                                                try:
                                                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                                        tmp_img.write(image_bytes)
                                                        tmp_img.flush()
                                                        uploaded_img = genai.upload_file(
                                                            path=tmp_img.name,
                                                            mime_type="image/png",
                                                            display_name=f"qa_image_{page_title}_{idx}.png"
                                                        )

                                                    full_prompt = (
                                                        "You're analyzing a technical image extracted from documentation. "
                                                        "Answer the user's question based on the visual content of the image, "
                                                        "as well as the summary below.\n\n"
                                                        f"Summary:\n{summary}\n\n"
                                                        f"User Question:\n{user_question}"
                                                    )

                                                    ai_response = ai_model.generate_content([uploaded_img, full_prompt])
                                                    st.session_state[ai_response_key] = ai_response.text.strip()
                                                except Exception as e:
                                                    st.error(f"AI failed to generate a response: {e}")

                                        # ✅ Display response if available
                                        if st.session_state.get(ai_response_key):
                                            st.info(st.session_state[ai_response_key])

                                        summary = st.session_state[summary_key]

                                        # 🔽 Then the download section
                                        file_name = st.text_input("Enter file name (no extension)", value=f"summary_image_{idx}", key=f"file_name_{page_title}_{idx}")
                                        format_choice = st.selectbox("Select format", ["PDF", "DOCX", "TXT", "Markdown"], key=f"format_{page_title}_{idx}")

                                        content, mime, ext = None, "", ""
                                        if format_choice == "PDF":
                                            content = generate_pdf(image_bytes, summary)
                                            mime = "application/pdf"
                                            ext = "pdf"
                                        elif format_choice == "DOCX":
                                            content = generate_docx(image_bytes, summary)
                                            mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                            ext = "docx"
                                        elif format_choice == "TXT":
                                            content = generate_txt(summary)
                                            mime = "text/plain"
                                            ext = "txt"
                                        elif format_choice == "Markdown":
                                            content = generate_md(image_bytes, summary)
                                            mime = "text/markdown"
                                            ext = "md"

                                        if content:
                                            st.download_button(
                                                label="Download",
                                                data=content,
                                                file_name=f"{file_name}.{ext}",
                                                mime=mime,
                                                key=f"download_button_{page_title}_{idx}"
                                            )

                                            

                                        # ✅ Create Graph button appears only after summary is ready
                                        if st.button("Create Graph", key=f"create_graph_{page_title}_{idx}"):
                                            with st.spinner(" Extracting data from image ..."):
                                                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                                    tmp_img.write(image_bytes)
                                                    tmp_img.flush()
                                                    uploaded_img = genai.upload_file(
                                                        path=tmp_img.name,
                                                        mime_type="image/png",
                                                        display_name=f"chart_image_{page_title}_{idx}.png"
                                                    )

                                                graph_prompt = (
                                                    "You're looking at a Likert-style bar chart image or table. Extract the full numeric table represented by the chart.\n"
                                                    "Return only the raw CSV table: no markdown, no comments, no code blocks.\n"
                                                    "The first column must be the response category (e.g., Strongly Agree), followed by columns for group counts (e.g., Students, Lecturers, Staff, Total).\n"
                                                    "Ensure all values are numeric and the CSV is properly aligned. Do NOT summarize—just output the table."
                                                )

                                                graph_response = ai_model.generate_content([uploaded_img, graph_prompt])
                                                csv_text = graph_response.text.strip()
                                                

                                                try:
                                                    cleaned_csv = clean_ai_csv(csv_text)
                                                    df = pd.read_csv(StringIO(cleaned_csv))

                                                    for col in df.columns[1:]:
                                                        df[col] = pd.to_numeric(df[col], errors='coerce')

                                                    df.dropna(subset=df.columns[1:], how='all', inplace=True)

                                                    if df.empty:
                                                        raise ValueError("Extracted DataFrame is empty after cleaning.")

                                                    st.session_state[graph_df_key] = df
                                                    st.subheader("Table from Image")
                                                    st.dataframe(df)
                                                    

                                                except Exception as e:
                                                    st.error(f"⚠️ Failed to parse chart data. Reason: {e}")
                                                    st.text_area("🧾 AI Output ", csv_text, height=200)

                                    if graph_df_key in st.session_state:
                                        df = st.session_state[graph_df_key]
                                        chart_type = st.selectbox("📈 Choose Chart Type", ["Grouped Bar", "Stacked Bar", "Line", "Pie"], key=chart_type_key)

                                        if chart_type == "Grouped Bar":
                                            fig = plot_grouped_bar(df)
                                        elif chart_type == "Stacked Bar":
                                            fig = plot_stacked_bar(df)
                                        elif chart_type == "Line":
                                            fig = plot_line(df)
                                        else:
                                            fig = plot_pie(df)

                                        st.session_state[graph_fig_key] = fig
                                        st.pyplot(fig)
                                        
                                        

                                        download_format = st.selectbox(
                                            "Download Format",
                                            ["PNG", "JPG", "SVG", "PDF", "DOCX", "PPTX"],
                                            key=f"download_fmt_{page_title}_{idx}"
                                        )

                                        chart_file_name = st.text_input("Enter chart file name", f"chart_{page_title}_{idx}", key=f"chart_file_name_{page_title}_{idx}")

                                        # Default values
                                        fmt_ext, mime, download_data = None, None, None

                                        if download_format == "PNG":
                                            fmt_ext, mime = "png", "image/png"
                                            download_data = get_image_bytes(fig, "png")
                                        elif download_format == "JPG":
                                            fmt_ext, mime = "jpg", "image/jpeg"
                                            download_data = get_image_bytes(fig, "jpg")
                                        elif download_format == "SVG":
                                            fmt_ext, mime = "svg", "image/svg+xml"
                                            download_data = get_image_bytes(fig, "svg")
                                        elif download_format == "PDF":
                                            fmt_ext, mime = "pdf", "application/pdf"
                                            download_data = get_image_bytes(fig, "pdf")
                                        elif download_format == "DOCX":
                                            from docx import Document
                                            from docx.shared import Inches

                                            output = BytesIO()
                                            doc = Document()
                                            image_stream = get_image_bytes(fig, "png")
                                            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                                tmp_img.write(image_stream.read())
                                                tmp_img.flush()
                                                doc.add_picture(tmp_img.name, width=Inches(6))
                                            doc.save(output)
                                            output.seek(0)
                                            fmt_ext, mime, download_data = "docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", output
                                        elif download_format == "PPTX":
                                            from pptx import Presentation
                                            from pptx.util import Inches

                                            prs = Presentation()
                                            slide_layout = prs.slide_layouts[5]  # blank layout
                                            slide = prs.slides.add_slide(slide_layout)

                                            image_stream = get_image_bytes(fig, "png")
                                            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                                tmp_img.write(image_stream.read())
                                                tmp_img.flush()
                                                slide.shapes.add_picture(tmp_img.name, Inches(1), Inches(1), Inches(6), Inches(4.5))

                                            output = BytesIO()
                                            prs.save(output)
                                            output.seek(0)
                                            fmt_ext, mime, download_data = "pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", output

                                        if download_data:
                                            st.download_button(
                                                label="Download Graph",
                                                data=download_data,
                                                file_name=f"{chart_file_name}.{fmt_ext}",
                                                mime=mime,
                                                key=f"download_chart_{page_title}_{idx}"
                                            )

                                        
                                except Exception as e:
                                    st.error(f"⚠️ Error processing image {idx}: {str(e)}")
            except Exception as e:
                st.error(f"❌ Error fetching pages or processing images: {str(e)}")
    else:
        st.error("❌ Failed to connect to Confluence.")

# ------------- Main App Dropdown -------------
feature_options = [
    "AI Powered Search",
    "Video Summarizer",
    "Code Assistant",
    "Impact Analyzer",
    "Test Support Tool",
    "Image"
]
# Create three columns, center column is wider
col1, col2, col3 = st.columns([1, 2, 1])
with col2:
    st.markdown(
        """
        <div style='text-align: center; margin-bottom: -30px;'>
            <span style='font-size: 2rem; font-weight: 800;'>Select a Feature</span>
        </div>
        """,
        unsafe_allow_html=True
    )
    selected_feature = st.selectbox(
        "", feature_options, index=None, placeholder="Select"
    )

if selected_feature == "AI Powered Search":
    feature_1()
elif selected_feature == "Video Summarizer":
    feature_2()
elif selected_feature == "Code Assistant":
    feature_3()
elif selected_feature == "Impact Analyzer":
    feature_4()
elif selected_feature == "Test Support Tool":
    feature_5()
elif selected_feature == "Image":
    feature_6()
# If nothing is selected, do nothing 
